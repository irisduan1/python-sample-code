from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from pptx.enum.chart import XL_LABEL_POSITION, XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from zipfile import ZipFile
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import shutil
import os
import collections
# import plotly.graph_objects as go
# import plotly.express as px
# from plotly.subplots import make_subplots
# from email.mime.multipart import MIMEMultipart
# from email.mime.text import MIMEText
# from email.mime.application import MIMEApplication
# from email.utils import COMMASPACE, formatdate

from os.path import basename
from datetime import datetime, timedelta,date
from dateutil.relativedelta import relativedelta
import calendar


## reading location and output location
folder = "/your-directory"
file_location = "/your-directory"
my_ppt = folder + "/report/template.pptx"
out_ppt = folder + "/report_{dt.year}{dt.month}{dt.day}.pptx".format(dt=datetime.now())


## read external files
numOfLines = pd.read_csv(folder+"/csv1.csv",index_col=0,skiprows=[2],header=None)


## date manipulate so that the program will auto match with current month to take current month number 
month_list_all = [date.today()+relativedelta(months=-i) for i in range(1,15)]
current_datetime = ['Yr' + str(i.year) + str(i.strftime('%B')) for i in month_list_all]
date_match = {key:value for key, value in zip(current_datetime,month_list_all)}
numOfLines.columns = numOfLines.iloc[0]+numOfLines.iloc[1]
numOfLines = numOfLines.loc['Line',current_datetime].to_frame().reset_index()
numOfLines.columns = ['date_long','lines']
numOfLines['Date'] = numOfLines.date_long.apply(lambda x: str(date_match.get(x))[:7])
numOfLines['lines'] = numOfLines.lines.apply(lambda x: int(x.strip().replace(',','')))
total_numOfLines=numOfLines[['Date','lines']]


    
prs = Presentation(my_ppt)

def slide_line_chart():
        
    shape_objects = prs.slides[0].shapes
    
    ## Read CSV, market level
    numOfUser = pd.read_csv(file_location + "csv2.csv",header=None)    
    numOfUser.columns = ['Date','Market','Count']   
    numOfUser.Count = round(numOfUser.Count/1000,0)
    numOfUser = numOfUser.groupby(['Date','Market']).sum().reset_index()
 
    ## National level data
    national_numOfUser = pd.read_csv(file_location + "csv3.csv",header=None)   
    national_numOfUser.columns = ['Date','Count']
    national_numOfUser.Count = round(national_numOfUser.Count/1000000,2)
    
    ## creating chart object where date is the x-axis
    chart_data = ChartData()
    chart_data.categories = [datetime.strptime(d,'%Y-%m').strftime('%Y-%m') for d in numOfUser.Date.unique()]
    
    ## one line per market
    for p in numOfUser.groupby(['Market']):
        
        chart_data.add_series('%s' % p[1].Market.unique()[0], p[1].Count.tolist())
    
    ## add national line to the chart
    chart_data.add_series('National', national_numOfUser.Count.tolist())
    
    x, y, cx, cy = Inches(0.5), Inches(1), Inches(12), Inches(5)
    chart = shape_objects.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    ## Left-side y axis
    y_axis = chart.value_axis
    y_axis.has_title = True 
    y_axis.axis_title.text_frame.text = 'Count(M)'
    y_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(13)
    y_axis.axis_title.text_frame.paragraphs[0].runs[0].font.bold= False
    y_axis.major_gridlines.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    y_axis.format.line.color.rgb = RGBColor(255, 255, 255)  
    y_axis.tick_labels.font.size = Pt(8) 
    
    ## X axis
    x_axis = chart.category_axis
    x_axis.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        
    ## Adjust line color, each for each market
    plot = chart.plots[0]
    plot.series[0].format.line.color.rgb = RGBColor(102, 51, 0)
    plot.series[1].format.line.color.rgb = RGBColor(225, 190, 15)
    plot.series[2].format.line.color.rgb = RGBColor(204, 229, 255)
    plot.series[3].format.line.color.rgb = RGBColor(32, 32, 32)
    plot.series[4].format.line.color.rgb = RGBColor(225, 169, 40)
    plot.series[5].format.line.color.rgb = RGBColor(0, 0, 204)
    plot.series[6].format.line.color.rgb = RGBColor(204, 255, 204)
    plot.series[7].format.line.color.rgb = RGBColor(255, 0, 0)
    plot.series[8].format.line.color.rgb = RGBColor(255, 153, 204) 
    plot.series[9].format.line.color.rgb = RGBColor(102, 178, 255) 
    plot.series[10].format.line.color.rgb = RGBColor(255, 153, 51)  
    plot.series[11].format.line.color.rgb = RGBColor(0, 76, 153)     
    
    plot.series[3].format.line.dash_style = MSO_LINE.DASH
    plot.series[4].format.line.dash_style = MSO_LINE.DASH
    plot.series[9].format.line.dash_style = MSO_LINE.DASH
    
    ## add line marker to national line so that easier to separate from markets

    plot.series[11].marker.style = XL_MARKER_STYLE.SQUARE
    plot.series[11].marker.size = 10

    marker_color = plot.series[11].marker.format
    marker_color.fill.solid()
    marker_color.fill.fore_color.rgb = RGBColor(0, 76, 153)
    marker_color.line.color.rgb = RGBColor(0, 76, 153)

        
    ## Add annotations to most top and bottom lines 
    value_series = np.asarray([list(l.values) for l in plot.series])
    value_series_max = dict(collections.Counter(value_series.argmax(axis=0)))
    value_series_min = dict(collections.Counter(value_series.argmin(axis=0)))
    max_series = max(value_series_max, key=value_series_max.get)
    min_series = max(value_series_min, key=value_series_min.get)
    
    for ind in range(len(chart_data.categories)):
        data_label_max = plot.series[max_series].points[ind].data_label
        data_label_max.font.size = Pt(8)
        # data_label_max.number_format = '0.00'
        data_label_max.position = XL_LABEL_POSITION.ABOVE 
        
        data_labels_min = plot.series[min_series].points[ind].data_label    
        data_labels_min.font.size = Pt(8)
        # data_labels_min.number_format = '0.00'
        data_labels_min.position = XL_LABEL_POSITION.BELOW  
    
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.size = Pt(9.5)
    chart.legend.include_in_layout = False
    chart.font.size = Pt(12)



if __name__ == '__main__':
    
    slide_line_chart()


 ## save and output the report, if the folder not exist - create the folder and save to it; if exist - replace the file inside the folder to most recent   
    try:
        shutil.rmtree(folder+'/report')
        os.mkdir(folder+'/report')
        prs.save(out_ppt)
    except:
        os.mkdir(folder+'/report')
        prs.save(out_ppt)



    
